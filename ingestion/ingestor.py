import os
import re
import io
import json
import hashlib
import logging
from pathlib import Path
from datetime import datetime

import fitz
import numpy as np
from PIL import Image

from app.core.config import (
    VAULT_PATH,
    ASSETS_PATH,
    SUPPORTED_TEXT,
    SUPPORTED_IMAGES,
    SUPPORTED_DOCS,
    OCR_LANGUAGES,
)

logger = logging.getLogger("Ingestor")


class Ingestor:

    # =====================================================
    # INIT
    # =====================================================

    def __init__(
        self,
        vault_path=VAULT_PATH,
    ):
        self.vault_path = Path(
            vault_path
        )

        self.assets_path = Path(
            self.vault_path / "assets"
        )
        self.vault_path.mkdir(parents=True, exist_ok=True)
        self.assets_path.mkdir(parents=True, exist_ok=True)

        self.reader = None

        self.saved_hashes = set()

    # =====================================================
    # OCR
    # =====================================================

    def get_ocr_reader(self):

        if self.reader is None:

            import easyocr

            logger.info(
                "Loading EasyOCR..."
            )

            self.reader = (
                easyocr.Reader(
                    OCR_LANGUAGES
                )
            )

        return self.reader

    def ocr_image(
        self,
        image_bytes,
    ):
        try:

            reader = (
                self.get_ocr_reader()
            )

            results = reader.readtext(
                image_bytes,
                detail=0,
            )

            return "\n".join(
                results
            )

        except Exception as e:

            logger.warning(
                f"OCR failed: {e}"
            )

            return ""

    # =====================================================
    # HELPERS
    # =====================================================

    def sanitize_title(
        self,
        title,
    ):
        return re.sub(
            r'[\\/*?:"<>|]',
            "",
            title,
        ).strip()

    def compute_hash(
        self,
        data,
    ):
        return hashlib.md5(
            data
        ).hexdigest()

    def timestamp(self):
        return datetime.now().strftime(
            "%Y-%m-%d %H:%M:%S"
        )

    # =====================================================
    # MARKDOWN
    # =====================================================

    def build_metadata(
        self,
        source,
        file_type,
    ):
        return (
            "---\n"
            f"source: {source}\n"
            f"type: {file_type}\n"
            f"date: {self.timestamp()}\n"
            "---\n\n"
        )

    def save_markdown(
        self,
        title,
        content,
        source,
        file_type,
    ):
        safe_title = (
            self.sanitize_title(
                title
            )
        )

        path = self.vault_path / f"{safe_title}.md"
        if path.exists():
            existing = path.read_text(encoding="utf-8", errors="ignore")
            if f"source: {source}\n" not in existing or f"type: {file_type}\n" not in existing:
                base_title = self.sanitize_title(f"{safe_title} ({file_type})")
                safe_title = base_title
                path = self.vault_path / f"{safe_title}.md"
                counter = 2
                while path.exists():
                    safe_title = f"{base_title} {counter}"
                    path = self.vault_path / f"{safe_title}.md"
                    counter += 1

        metadata = (
            self.build_metadata(
                source,
                file_type,
            )
        )

        with open(
            path,
            "w",
            encoding="utf-8",
        ) as f:

            f.write(
                metadata
                + content
            )

        logger.info(
            f"Saved {safe_title}"
        )

        return safe_title

    # =====================================================
    # DISPATCHER
    # =====================================================

    def ingest(
        self,
        file_path,
    ):
        file_path = Path(
            file_path
        )

        if not file_path.exists():

            raise FileNotFoundError(
                file_path
            )

        ext = (
            file_path.suffix
            .lower()
        )

        if ext in SUPPORTED_TEXT:

            return (
                self.ingest_text(
                    file_path
                )
            )

        if ext in SUPPORTED_IMAGES:

            return (
                self.ingest_image(
                    file_path
                )
            )

        if ext == ".pdf":
            return self.ingest_pdf(file_path)
        if ext == ".docx":
            return self.ingest_docx(file_path)
        if ext == ".pptx":
            return self.ingest_pptx(file_path)

        raise ValueError(
            f"Unsupported file: {ext}"
        )

    # =====================================================
    # TEXT INGESTION
    # =====================================================

    def ingest_text(
        self,
        file_path,
    ):
        title = (
            file_path.stem
        )

        with open(
            file_path,
            "r",
            encoding="utf-8",
            errors="ignore",
        ) as f:

            content = f.read()

        markdown = (
            f"# {title}\n\n"
            "#pending\n\n"
            f"{content}\n"
        )

        return self.save_markdown(
            title,
            markdown,
            str(file_path),
            "text",
        )

    # =====================================================
    # IMAGE INGESTION
    # =====================================================

    def save_asset(
        self,
        image_bytes,
        extension=".png",
    ):
        image_hash = (
            self.compute_hash(
                image_bytes
            )
        )

        if (
            image_hash
            in self.saved_hashes
        ):
            return None

        self.saved_hashes.add(
            image_hash
        )

        filename = (
            f"{image_hash}"
            f"{extension}"
        )

        path = (
            self.assets_path
            / filename
        )

        with open(
            path,
            "wb",
        ) as f:

            f.write(
                image_bytes
            )

        return filename

    def ingest_image(
        self,
        file_path,
    ):
        title = (
            file_path.stem
        )

        with open(
            file_path,
            "rb",
        ) as f:

            image_bytes = (
                f.read()
            )

        extension = (
            file_path.suffix
        )

        asset_name = (
            self.save_asset(
                image_bytes,
                extension,
            )
        )

        ocr_text = (
            self.ocr_image(
                image_bytes
            )
        )

        markdown = (
            f"# {title}\n\n"
            "#pending\n\n"
        )

        if asset_name:

            markdown += (
                f"![[{asset_name}]]\n\n"
            )

        markdown += (
            "## OCR Text\n\n"
            f"{ocr_text}\n"
        )

        return self.save_markdown(
            title,
            markdown,
            str(file_path),
            "image",
        )

    # =====================================================
    # OFFICE DOCUMENTS
    # =====================================================

    def ingest_docx(self, file_path):
        from docx import Document
        document = Document(str(file_path))
        lines = [f"# {file_path.stem}\n", "#pending\n"]
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style = paragraph.style.name.lower()
            lines.append(f"{'## ' if 'heading' in style else ''}{text}\n")
        for table in document.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if rows:
                lines.extend(["\n## Table\n", "| " + " | ".join(rows[0]) + " |\n",
                              "| " + " | ".join("---" for _ in rows[0]) + " |\n"])
                lines.extend("| " + " | ".join(row) + " |\n" for row in rows[1:])
        return self.save_markdown(file_path.stem, "\n".join(lines), str(file_path), "docx")

    def ingest_pptx(self, file_path):
        from pptx import Presentation
        presentation = Presentation(str(file_path))
        lines = [f"# {file_path.stem}\n", "#pending\n"]
        for index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"\n## Slide {index}\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip() + "\n")
        return self.save_markdown(file_path.stem, "\n".join(lines), str(file_path), "pptx")

    # =====================================================
    # IMAGE UTILITIES
    # =====================================================

    def pil_to_bytes(
        self,
        image,
    ):
        buffer = io.BytesIO()

        image.save(
            buffer,
            format="PNG",
        )

        return (
            buffer.getvalue()
        )

    def extract_image_text(
        self,
        image,
    ):
        image_bytes = (
            self.pil_to_bytes(
                image
            )
        )

        return self.ocr_image(
            image_bytes
        )

    # =====================================================
    # PDF PLACEHOLDER
    # =====================================================

    
        # =====================================================
    # PDF HELPERS
    # =====================================================

    def get_header_thresholds(
        self,
        doc,
    ):
        sizes = []

        sample_pages = min(
            len(doc),
            10,
        )

        for page_num in range(
            sample_pages
        ):

            page = doc[
                page_num
            ]

            blocks = (
                page.get_text(
                    "dict"
                )["blocks"]
            )

            for block in blocks:

                if (
                    "lines"
                    not in block
                ):
                    continue

                for line in block[
                    "lines"
                ]:

                    for span in line[
                        "spans"
                    ]:

                        sizes.append(
                            span[
                                "size"
                            ]
                        )

        if not sizes:

            return (
                18,
                15,
            )

        h1 = np.percentile(
            sizes,
            99.5,
        )

        h2 = np.percentile(
            sizes,
            98,
        )

        return h1, h2


    def classify_header(
        self,
        text,
        size,
        h1,
        h2,
    ):
        text = text.strip()

        if not text:
            return 0

        if (
            size >= h1
            and len(text) < 120
        ):
            return 1

        if (
            size >= h2
            and len(text) < 120
        ):
            return 2

        return 0


    def extract_text_block(
        self,
        block,
    ):
        text = ""

        max_size = 0

        for line in block[
            "lines"
        ]:

            for span in line[
                "spans"
            ]:

                span_text = (
                    span["text"]
                    .replace(
                        "\x00",
                        ""
                    )
                    .strip()
                )

                if not span_text:
                    continue

                text += (
                    span_text
                    + " "
                )

                max_size = max(
                    max_size,
                    span["size"],
                )

        return (
            text.strip(),
            max_size,
        )

    # =====================================================
    # IMAGE EXTRACTION
    # =====================================================

    def extract_page_images(
        self,
        doc,
        page,
    ):
        image_refs = []

        for image in page.get_images(
            full=True
        ):

            try:

                xref = image[0]

                base_image = (
                    doc.extract_image(
                        xref
                    )
                )

                image_bytes = (
                    base_image[
                        "image"
                    ]
                )

                ext = (
                    "."
                    + base_image.get(
                        "ext",
                        "png",
                    )
                )

                asset_name = (
                    self.save_asset(
                        image_bytes,
                        ext,
                    )
                )

                if asset_name:

                    image_refs.append(
                        asset_name
                    )

            except Exception:

                continue

        return image_refs

    # =====================================================
    # TABLE EXTRACTION
    # =====================================================

    def extract_tables(
        self,
        page,
    ):
        tables_md = []

        try:

            tables = (
                page.find_tables()
            )

            for table in tables:

                try:

                    df = (
                        table.to_pandas()
                    )

                    md = (
                        df.to_markdown(
                            index=False
                        )
                    )

                    tables_md.append(
                        md
                    )

                except Exception:
                    pass

        except Exception:
            pass

        return tables_md

    # =====================================================
    # OCR FALLBACK
    # =====================================================

    def ocr_pdf_page(
        self,
        page,
    ):
        pix = page.get_pixmap(
            matrix=fitz.Matrix(
                2,
                2,
            )
        )

        image_bytes = pix.tobytes(
            "png"
        )

        return self.ocr_image(
            image_bytes
        )

    # =====================================================
    # PDF INGESTION
    # =====================================================

    def ingest_pdf(
        self,
        file_path,
    ):
        title = (
            file_path.stem
        )

        doc = fitz.open(
            str(file_path)
        )

        h1_thresh, h2_thresh = (
            self.get_header_thresholds(
                doc
            )
        )

        markdown = []

        markdown.append(
            f"# {title}\n"
        )

        markdown.append(
            "#pending\n"
        )

        markdown.append(
            "## 📖 Study Outline\n"
        )

        toc_entries = []

        all_content = []

        for page_number in range(
            len(doc)
        ):

            page = doc[
                page_number
            ]

            page_content = []

            page_content.append(
                f"\n---\n"
                f"Page {page_number+1}\n"
                f"---\n"
            )

            # -------------------
            # Tables
            # -------------------

            tables = (
                self.extract_tables(
                    page
                )
            )

            for table_md in tables:

                page_content.append(
                    "\n## Table\n"
                )

                page_content.append(
                    table_md
                )

            # -------------------
            # Text
            # -------------------

            text_found = False

            blocks = (
                page.get_text(
                    "dict"
                )["blocks"]
            )

            for block in blocks:

                if (
                    block.get(
                        "type",
                        0
                    )
                    != 0
                ):
                    continue

                if (
                    "lines"
                    not in block
                ):
                    continue

                text, size = (
                    self.extract_text_block(
                        block
                    )
                )

                if not text:
                    continue

                text_found = True

                level = (
                    self.classify_header(
                        text,
                        size,
                        h1_thresh,
                        h2_thresh,
                    )
                )

                if level > 0:

                    header = (
                        "#"
                        * (
                            level + 1
                        )
                    )

                    page_content.append(
                        f"\n{header} {text}\n"
                    )

                    toc_entries.append(
                        text
                    )

                else:

                    page_content.append(
                        text + "\n"
                    )

            # -------------------
            # OCR fallback
            # -------------------

            if not text_found:

                ocr_text = (
                    self.ocr_pdf_page(
                        page
                    )
                )

                if ocr_text:

                    page_content.append(
                        "\n## OCR Content\n"
                    )

                    page_content.append(
                        ocr_text
                    )

            # -------------------
            # Images
            # -------------------

            image_refs = (
                self.extract_page_images(
                    doc,
                    page,
                )
            )

            for img in image_refs:

                page_content.append(
                    f"\n![[{img}]]\n"
                )

            all_content.extend(
                page_content
            )

        # -------------------
        # TOC
        # -------------------

        for entry in toc_entries[
            :100
        ]:

            markdown.append(
                f"- {entry}\n"
            )

        markdown.append(
            "\n---\n"
        )

        markdown.extend(
            all_content
        )

        doc.close()

        final_markdown = (
            "\n".join(
                markdown
            )
        )

        return self.save_markdown(
            title,
            final_markdown,
            str(file_path),
            "pdf",
        )
